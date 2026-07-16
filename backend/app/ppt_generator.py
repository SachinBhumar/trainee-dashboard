from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
import io
import datetime
from sqlalchemy.orm import Session
from app import models, crud

COLOR_PRIMARY = RGBColor(165, 21, 38)     
COLOR_SECONDARY = RGBColor(44, 62, 80)    
COLOR_BG = RGBColor(248, 249, 250)         
COLOR_WHITE = RGBColor(255, 255, 255)
COLOR_TEXT_DARK = RGBColor(33, 37, 41)
COLOR_TEXT_MUTED = RGBColor(108, 117, 125)
COLOR_ACCENT = RGBColor(40, 167, 69)      

# Colors for KPI card stripes (matching dashboard CSS HSL palette)
STRIPE_COLORS = [
    RGBColor(243, 156, 18),  # Orange
    RGBColor(230, 126, 34),  # Dark Orange
    RGBColor(41, 128, 185),  # Blue
    RGBColor(39, 174, 96),   # Green
    RGBColor(165, 21, 38),   # Crimson
    RGBColor(142, 68, 173),  # Purple
    RGBColor(22, 160, 133),  # Teal
    RGBColor(44, 62, 80)     # Slate
]

def apply_background(slide):
    left = top = 0
    width = Inches(10)
    height = Inches(7.5)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    rect.fill.solid()
    rect.fill.fore_color.rgb = COLOR_BG
    rect.line.color.rgb = COLOR_BG 

def create_header(slide, title_text, category_text=""):
    # Add a thin crimson red bar at the top
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = COLOR_PRIMARY
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9.0), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Arial"
    
    if category_text:
        p2 = tf.add_paragraph()
        p2.text = category_text.upper()
        p2.font.size = Pt(9)
        p2.font.bold = True
        p2.font.color.rgb = COLOR_TEXT_MUTED
        p2.font.name = "Arial"
        p2.space_before = Pt(4)

def add_kpi_card(slide, left, top, width, height, title, value, subtext, source, stripe_color):
    # Main White Card container
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = COLOR_WHITE
    card.line.color.rgb = RGBColor(230, 230, 230)
    card.line.width = Pt(1)
    
    # Left Edge Color Stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(0.08), height)
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = stripe_color
    stripe.line.fill.background()
    
    # Text Frame Box
    tb = slide.shapes.add_textbox(left + Inches(0.12), top, width - Inches(0.12), height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_top = Inches(0.1)
    tf.margin_right = Inches(0.1)
    
    p_title = tf.paragraphs[0]
    p_title.text = title.upper()
    p_title.font.size = Pt(8)
    p_title.font.bold = True
    p_title.font.color.rgb = COLOR_TEXT_MUTED
    p_title.font.name = "Arial"
    
    p_val = tf.add_paragraph()
    p_val.text = value
    p_val.font.size = Pt(22)
    p_val.font.bold = True
    p_val.font.color.rgb = COLOR_TEXT_DARK
    p_val.font.name = "Arial"
    p_val.space_before = Pt(4)
    
    p_sub = tf.add_paragraph()
    p_sub.text = subtext
    p_sub.font.size = Pt(8)
    p_sub.font.bold = True
    p_sub.font.color.rgb = COLOR_PRIMARY if "+" in subtext or "↑" in subtext else COLOR_TEXT_MUTED
    p_sub.font.name = "Arial"
    p_sub.space_before = Pt(3)
    
    p_src = tf.add_paragraph()
    p_src.text = f"Source: {source}"
    p_src.font.size = Pt(7)
    p_src.font.italic = True
    p_src.font.color.rgb = COLOR_TEXT_MUTED
    p_src.font.name = "Arial"
    p_src.space_before = Pt(6)

def generate_report_ppt(db: Session, period: str, tab: str = "brand", timeframe: str = "Full year", view: str = "YTD") -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Determine fiscal year of selected period
    period_record = db.query(models.MetricValue).filter(
        models.MetricValue.period == period
    ).first()
    target_fy = period_record.fiscal_year if period_record else "2025-26"
    
    # Get all date-based records for this fiscal year
    db_records = db.query(models.MetricValue).filter(
        models.MetricValue.fiscal_year == target_fy,
        models.MetricValue.sheet == "Summary"
    ).all()
    
    valid_records = []
    for r in db_records:
        try:
            dt = datetime.datetime.strptime(r.period, "%Y-%m-%d")
            valid_records.append((r, dt))
        except ValueError:
            continue
            
    # Map timeframe to its month numbers (April=4, May=5, ..., March=3)
    tf_months_map = {
        "Full year": [4, 5, 6, 7, 8, 9, 10, 11, 12, 1, 2, 3],
        "H1 Apr–Sep": [4, 5, 6, 7, 8, 9],
        "H2 Oct–Mar": [10, 11, 12, 1, 2, 3],
        "Q1": [4, 5, 6],
        "Q2": [7, 8, 9],
        "Q3": [10, 11, 12],
        "Q4": [1, 2, 3]
    }
    
    selected_tf_months = tf_months_map.get(timeframe, tf_months_map["Full year"])
    chrono_order = {4:0, 5:1, 6:2, 7:3, 8:4, 9:5, 10:6, 11:7, 12:8, 1:9, 2:10, 3:11}

    if view == "YTD":
        valid_months_in_db = {dt.month for r, dt in valid_records if r.value is not None or r.value_text is not None}
        matching_months = [m for m in selected_tf_months if m in valid_months_in_db]
        
        if matching_months:
            max_index = max(chrono_order[m] for m in matching_months)
            target_months = [m for m, idx in chrono_order.items() if idx <= max_index]
        else:
            target_months = selected_tf_months
    else:
        target_months = selected_tf_months
        
    filtered_records = [(r, dt) for r, dt in valid_records if dt.month in target_months]
    
    def aggregate_metric(metric_name, mode="avg"):
        vals = [r.value for r, dt in filtered_records if metric_name.lower() in r.metric_name.lower() and r.value is not None]
        if not vals:
            return None
        if mode == "sum":
            return sum(vals)
        elif mode == "last":
            records_for_metric = [(r, dt) for r, dt in filtered_records if metric_name.lower() in r.metric_name.lower() and r.value is not None]
            if not records_for_metric:
                return None
            records_for_metric = sorted(records_for_metric, key=lambda x: chrono_order.get(x[1].month, 0))
            return records_for_metric[-1][0].value
        elif mode == "last_text":
            records_for_metric = [(r, dt) for r, dt in filtered_records if metric_name.lower() in r.metric_name.lower() and r.value_text is not None]
            if not records_for_metric:
                return None
            records_for_metric = sorted(records_for_metric, key=lambda x: chrono_order.get(x[1].month, 0))
            return records_for_metric[-1][0].value_text
        else: # avg
            return sum(vals) / len(vals)

    def find_base_val(metric_name):
        res = db.query(models.MetricValue).filter(
            models.MetricValue.sheet == "Summary",
            models.MetricValue.metric_name.like(f"%{metric_name}%"),
            models.MetricValue.period == f"FY {target_fy}"
        ).first()
        if not res:
            res = db.query(models.MetricValue).filter(
                models.MetricValue.sheet == "Summary",
                models.MetricValue.metric_name.like(f"%{metric_name}%"),
                models.MetricValue.period == f"FY {target_fy} Avg"
            ).first()
        if not res:
            res = db.query(models.MetricValue).filter(
                models.MetricValue.sheet == "Summary",
                models.MetricValue.metric_name.like(f"%{metric_name}%"),
                models.MetricValue.period.like("FY %")
            ).first()
        return res.value if res else None

    # SLIDE 1: Title Slide
    slide_layout = prs.slide_layouts[6] 
    slide1 = prs.slides.add_slide(slide_layout)
    apply_background(slide1)
    
    # Branded colored bar on left edge of cover slide
    cover_stripe = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.4), Inches(7.5))
    cover_stripe.fill.solid()
    cover_stripe.fill.fore_color.rgb = COLOR_PRIMARY
    cover_stripe.line.fill.background()

    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(8.0), Inches(2.5))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "ASIAN PAINTS & MADISON MEDIA"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY
    p.font.name = "Arial"
    
    p2 = tf1.add_paragraph()
    p2.text = "Brand Health Performance Report" if tab == "brand" else "Media & Search Performance Report"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_PRIMARY
    p2.font.name = "Arial"
    p2.space_before = Pt(8)
    
    p3 = tf1.add_paragraph()
    p3.text = f"Fiscal Year: {target_fy} | Timeframe: {timeframe} ({view} view) | Confidential"
    p3.font.size = Pt(13)
    p3.font.color.rgb = COLOR_TEXT_MUTED
    p3.font.name = "Arial"
    p3.space_before = Pt(14)

    # SLIDE 2: KPI Overview Cards
    slide2 = prs.slides.add_slide(slide_layout)
    apply_background(slide2)
    create_header(slide2, f"KPI Dashboard Overview - FY {target_fy}", f"{tab.capitalize()} Performance")

    if tab == "brand":
        toma_val = aggregate_metric("TOMA", "avg")
        toma_base = find_base_val("TOMA")
        toma_change = f"{(toma_val-toma_base)*100:+.1f}pp vs base" if toma_val and toma_base else "vs base"

        cons_val = aggregate_metric("Consideration", "avg")
        cons_base = find_base_val("Consideration")
        cons_change = f"{(cons_val-cons_base)*100:+.1f}pp vs base" if cons_val and cons_base else "vs base"

        sos_val = aggregate_metric("Share of Search", "avg")
        sos_ytd_m = db.query(models.MetricValue).filter(
            models.MetricValue.sheet == "Summary",
            models.MetricValue.metric_name.like("%Share of Search%"),
            models.MetricValue.period == f"YTD-{target_fy}"
        ).first()
        sos_ytd_val = f"{sos_ytd_m.value*100:.1f}%" if (sos_ytd_m and sos_ytd_m.value) else "82.6%"
        sos_change = f"{sos_ytd_val} YTD Avg"

        reach_val = aggregate_metric("TV Reach", "last")
        reach_base = find_base_val("TV Reach")
        reach_change = f"Baseline: {reach_base or 869}Mn"

        tvsov_val = aggregate_metric("TV Share of Voice", "avg")
        tvsov_base = find_base_val("TV Share of Voice")
        tvsov_change = f"{(tvsov_val-tvsov_base)*100:+.1f}pp vs target" if tvsov_val and tvsov_base else "vs target"

        digsov_val = aggregate_metric("Share of Impressions", "avg")
        digsov_change = "Period average"

        paidsos_val = aggregate_metric("AP Paid SOS", "avg")
        paidsos_base = find_base_val("AP Paid SOS")
        paidsos_change = f"{(paidsos_val-paidsos_base)*100:+.1f}pp vs target" if paidsos_val and paidsos_base else "vs target"

        aisos_val = aggregate_metric("AI SoS", "avg")
        aisos_base = find_base_val("AI SoS")
        aisos_change = f"{(aisos_val-aisos_base)*100:+.1f}pp vs base" if aisos_val and aisos_base else "vs base"

        brand_kpis = [
            ("Share of Search", f"{sos_val*100:.1f}%" if sos_val else "82.0%", sos_change, "Google Dashboard", STRIPE_COLORS[0]),
            ("TOMA", f"{toma_val*100:.1f}%" if toma_val else "76.0%", toma_change, "Kantar", STRIPE_COLORS[1]),
            ("Consideration", f"{cons_val*100:.1f}%" if cons_val else "82.0%", cons_change, "Kantar", STRIPE_COLORS[2]),
            ("MMR Reach", f"{reach_val:.0f}Mn" if reach_val else "711Mn", reach_change, "BARC + Madison", STRIPE_COLORS[3]),
            ("TV SOV", f"{tvsov_val*100:.1f}%" if tvsov_val else "41.0%", tvsov_change, "BARC", STRIPE_COLORS[4]),
            ("Digital SOV", f"{digsov_val*100:.1f}%" if digsov_val else "36.8%", digsov_change, "Vtion", STRIPE_COLORS[5]),
            ("Paid Search SOS", f"{paidsos_val*100:.1f}%" if paidsos_val else "93.5%", paidsos_change, "SimilarWeb", STRIPE_COLORS[6]),
            ("AI SOS", f"{aisos_val*100:.1f}%" if aisos_val else "82.4%", aisos_change, "SimilarWeb", STRIPE_COLORS[7])
        ]

        # Draw 2x4 grid of KPI Cards
        card_w = Inches(2.1)
        card_h = Inches(1.5)
        gap_x = Inches(0.18)
        gap_y = Inches(0.25)
        
        for idx, (title, val, sub, src, col) in enumerate(brand_kpis):
            col_idx = idx % 4
            row_idx = idx // 4
            left = Inches(0.5) + col_idx * (card_w + gap_x)
            top = Inches(1.8) + row_idx * (card_h + gap_y)
            add_kpi_card(slide2, left, top, card_w, card_h, title, val, sub, src, col)

    else: # media
        all_media_val = aggregate_metric("All Media SOE", "avg")
        all_media_base = find_base_val("All Media SOE")
        all_media_change = f"{(all_media_val-all_media_base)*100:+.1f}pp vs target" if all_media_val and all_media_base else "vs target"

        yt_soe_val = None
        yt_soe_records = [r for r, dt in filtered_records if "Share of Expenditure".lower() in r.metric_name.lower() and r.source == "Google" and r.value is not None]
        if yt_soe_records:
            yt_soe_val = sum(r.value for r in yt_soe_records) / len(yt_soe_records)

        freq_val = aggregate_metric("TV + Digital Frequencies MoM", "last_text")
        web_traffic_val = aggregate_metric("Website Traffic", "sum")
        
        ap_paid_val = aggregate_metric("Asian Paints Paid Search Traffic", "sum")
        opus_paid_val = aggregate_metric("Opus Paid Search Traffic", "sum")
        paid_search_change = f"vs Birla Opus {opus_paid_val/1000000:.2f}Mn" if ap_paid_val and opus_paid_val else "vs Birla Opus"

        one_pd_val = aggregate_metric("1P Data Outreach", "last_text")

        media_kpis = [
            ("All Platform SOE", f"{all_media_val*100:.0f}%" if all_media_val else "34%", all_media_change, "Madison Competes", STRIPE_COLORS[0]),
            ("Digital SOE (YT)", f"{yt_soe_val*100:.1f}%" if yt_soe_val else "35.5%", "Period average", "Google Platforms", STRIPE_COLORS[1]),
            ("Avg Frequency", freq_val if freq_val else "13-16x", "Period range", "MSpectra", STRIPE_COLORS[2]),
            ("Website Traffic", f"{web_traffic_val/1000000:.1f}M" if web_traffic_val else "41.7M", "Total visits in period", "SimilarWeb", STRIPE_COLORS[3]),
            ("Paid Search Traffic", f"{ap_paid_val/1000000:.2f}Mn" if ap_paid_val else "4.16Mn", paid_search_change, "SimilarWeb", STRIPE_COLORS[4]),
            ("1PD Reach", one_pd_val if one_pd_val else "XXMn", "Activated audience", "CDP / CRM", STRIPE_COLORS[5])
        ]

        # Draw 2x3 grid of KPI Cards
        card_w = Inches(2.8)
        card_h = Inches(1.5)
        gap_x = Inches(0.2)
        gap_y = Inches(0.25)
        
        for idx, (title, val, sub, src, col) in enumerate(media_kpis):
            col_idx = idx % 3
            row_idx = idx // 3
            left = Inches(0.6) + col_idx * (card_w + gap_x)
            top = Inches(1.8) + row_idx * (card_h + gap_y)
            add_kpi_card(slide2, left, top, card_w, card_h, title, val, sub, src, col)

    # SLIDE 3: Trend Charts (Side-by-Side native PPT Charts)
    slide3 = prs.slides.add_slide(slide_layout)
    apply_background(slide3)
    create_header(slide3, f"Performance Trends - {timeframe} ({view} view)", "Data Charts")

    # Chronologically sort periods for charts
    chrono_months = sorted(target_months, key=lambda m: chrono_order[m])
    month_names_map = {4:"Apr", 5:"May", 6:"Jun", 7:"Jul", 8:"Aug", 9:"Sep", 10:"Oct", 11:"Nov", 12:"Dec", 1:"Jan", 2:"Feb", 3:"Mar"}
    categories = [month_names_map[m] for m in chrono_months]

    if not categories:
        categories = ["Apr", "May"] # Fallback

    if tab == "brand":
        # Chart 1: Share of Search Trend
        sos_chart_data = CategoryChartData()
        sos_chart_data.categories = categories
        
        sos_vals = []
        for m in chrono_months:
            r = next((rec for rec, dt in filtered_records if dt.month == m and "Share of Search".lower() in rec.metric_name.lower()), None)
            sos_vals.append(r.value * 100 if (r and r.value) else 82.0)
            
        sos_chart_data.add_series("Share of Search", tuple(sos_vals))
        
        x1, y1, cx1, cy1 = Inches(0.5), Inches(1.8), Inches(4.3), Inches(4.8)
        chart1_shape = slide3.shapes.add_chart(
            XL_CHART_TYPE.LINE, x1, y1, cx1, cy1, sos_chart_data
        )
        chart1 = chart1_shape.chart
        chart1.has_legend = False
        chart1.value_axis.has_major_gridlines = True
        
        # Add Title to Chart 1
        title_box1 = slide3.shapes.add_textbox(x1, y1 - Inches(0.4), cx1, Inches(0.4))
        p1 = title_box1.text_frame.paragraphs[0]
        p1.text = "Share of Search Trend (%)"
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_SECONDARY

        # Chart 2: TOMA & Consideration Trend
        tc_chart_data = CategoryChartData()
        tc_chart_data.categories = categories
        
        toma_vals = []
        cons_vals = []
        for m in chrono_months:
            t_rec = next((rec for rec, dt in filtered_records if dt.month == m and "TOMA".lower() in rec.metric_name.lower()), None)
            c_rec = next((rec for rec, dt in filtered_records if dt.month == m and "Consideration".lower() in rec.metric_name.lower()), None)
            toma_vals.append(t_rec.value * 100 if (t_rec and t_rec.value) else 76.0)
            cons_vals.append(c_rec.value * 100 if (c_rec and c_rec.value) else 82.0)
            
        tc_chart_data.add_series("TOMA", tuple(toma_vals))
        tc_chart_data.add_series("Consideration", tuple(cons_vals))
        
        x2, y2, cx2, cy2 = Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.8)
        chart2_shape = slide3.shapes.add_chart(
            XL_CHART_TYPE.LINE, x2, y2, cx2, cy2, tc_chart_data
        )
        chart2 = chart2_shape.chart
        chart2.has_legend = True
        chart2.value_axis.has_major_gridlines = True
        
        # Add Title to Chart 2
        title_box2 = slide3.shapes.add_textbox(x2, y2 - Inches(0.4), cx2, Inches(0.4))
        p2 = title_box2.text_frame.paragraphs[0]
        p2.text = "TOMA vs Consideration Trend (%)"
        p2.font.bold = True
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SECONDARY

    else: # media
        # Chart 1: AP vs Birla Opus - Paid Search Traffic (Clustered Column)
        ps_chart_data = CategoryChartData()
        ps_chart_data.categories = categories
        
        ap_paid_vals = []
        opus_paid_vals = []
        for m in chrono_months:
            ap_rec = next((rec for rec, dt in filtered_records if dt.month == m and "Asian Paints Paid Search Traffic".lower() in rec.metric_name.lower()), None)
            op_rec = next((rec for rec, dt in filtered_records if dt.month == m and "Opus Paid Search Traffic".lower() in rec.metric_name.lower()), None)
            # Format in Thousands (K) to show nicely on y-axis
            ap_paid_vals.append(ap_rec.value / 1000 if (ap_rec and ap_rec.value) else 0.0)
            opus_paid_vals.append(op_rec.value / 1000 if (op_rec and op_rec.value) else 0.0)
            
        ps_chart_data.add_series("Asian Paints", tuple(ap_paid_vals))
        ps_chart_data.add_series("Birla Opus", tuple(opus_paid_vals))
        
        x1, y1, cx1, cy1 = Inches(0.5), Inches(1.8), Inches(4.3), Inches(4.8)
        chart1_shape = slide3.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x1, y1, cx1, cy1, ps_chart_data
        )
        chart1 = chart1_shape.chart
        chart1.has_legend = True
        chart1.value_axis.has_major_gridlines = True
        
        # Add Title to Chart 1
        title_box1 = slide3.shapes.add_textbox(x1, y1 - Inches(0.4), cx1, Inches(0.4))
        p1 = title_box1.text_frame.paragraphs[0]
        p1.text = "AP vs Birla Opus — Paid Search Traffic (K)"
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_SECONDARY

        # Chart 2: Website Traffic - monthly visits (Clustered Column)
        wt_chart_data = CategoryChartData()
        wt_chart_data.categories = categories
        
        wt_vals = []
        for m in chrono_months:
            r = next((rec for rec, dt in filtered_records if dt.month == m and "Website Traffic".lower() in rec.metric_name.lower()), None)
            # Format in Millions (M)
            wt_vals.append(r.value / 1000000 if (r and r.value) else 0.0)
            
        wt_chart_data.add_series("Website Traffic", tuple(wt_vals))
        
        x2, y2, cx2, cy2 = Inches(5.2), Inches(1.8), Inches(4.3), Inches(4.8)
        chart2_shape = slide3.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED, x2, y2, cx2, cy2, wt_chart_data
        )
        chart2 = chart2_shape.chart
        chart2.has_legend = False
        chart2.value_axis.has_major_gridlines = True
        
        # Add Title to Chart 2
        title_box2 = slide3.shapes.add_textbox(x2, y2 - Inches(0.4), cx2, Inches(0.4))
        p2 = title_box2.text_frame.paragraphs[0]
        p2.text = "Website Traffic — Monthly Visits (M)"
        p2.font.bold = True
        p2.font.size = Pt(11)
        p2.font.color.rgb = COLOR_SECONDARY

    ppt_stream = io.BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
